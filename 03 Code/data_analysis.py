import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from statsmodels.stats.descriptivestats import Description
from statsmodels.tsa.arima.model import ARIMA
from statsmodels.tsa.statespace.sarimax import SARIMAX
import seaborn as sns
from numpy.linalg import LinAlgError
import numpy as np
from statsmodels.tsa.stattools import ccovf, ccf, grangercausalitytests, adfuller
from sklearn.preprocessing import StandardScaler, MinMaxScaler
from scipy.stats import pearsonr, spearmanr
from sklearn.feature_selection import mutual_info_regression
from dtaidistance import dtw
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from PIL import Image
import plotly.express as px
import kaleido



# Single Time Series Analysis

# Understanding
def print_uniques(df: pd.DataFrame, columns=None, excluded=None, max_uniques=1500):
    excluded = excluded or []
    for column in columns if columns else df.columns:
        if column not in excluded:
            unique_values = set(df[column])
            if len(unique_values) > max_uniques:
                print(f"Uniques of '{column}': ({len(unique_values)}) Too many to print. Data Type: {df[column].dtype}")
            else:
                print(f"Uniques of '{column}': ({len(unique_values)}) {unique_values}")
# Analysis
def calculate_moving_statistics(dates, values, window) -> pd.DataFrame: 
    """
    Calculates various moving statistics for a time series

    :param dates: List of dates (datetime objects or strings in the format 'YYYY-MM-DD')
    :param values: List of values (float or int)
    :param window: Window size for calculating the moving average and standard deviation
    :return: DataFrame with original values, moving average, moving standard deviation
    """

    # Convert dates to datetime objects if needed
    if isinstance(dates[0], str):
        dates = pd.to_datetime(dates)

    # Create a DataFrame
    data = pd.DataFrame({'Date': dates, 'Value': values})
    data.set_index('Date', inplace=True)

    # Calculate moving average, moving standard deviation, moving variance, moving median, moving range
    data['Moving Average'] = data['Value'].rolling(window=window).mean()
    data['Moving Std Dev'] = data['Value'].rolling(window=window).std()
    data['Moving Median'] = data['Value'].rolling(window=window).median()
    data['Moving Range'] = data['Value'].rolling(window=window).apply(lambda x: x.max() - x.min(), raw=True)

    return data
def calculate_arima_sarima(data, window, arima_order=(3, 1, 3), sarima_order=(1, 1, 2, 12)):
    """
    Calculates ARIMA and SARIMA predictions for a given time series DataFrame.

    :param data: DataFrame with a 'Value' column containing the time series data
    :param window: Window size for calculating the moving statistics
    :param arima_order: Order tuple (p,d,q) for the ARIMA model
    :param sarima_order: Order tuple (p,d,q,s) for the SARIMA model
    :return: DataFrame with added columns for ARIMA and SARIMA predictions
    """
    
    # Prepare to store predictions
    arima_predictions = []
    sarima_predictions = []

    # Ensure the DataFrame has a regular frequency
    data = data.asfreq('MS')  # Assuming monthly data, change 'MS' if needed

    # Loop through the data to generate predictions
    for i in range(window, len(data)):
        # Use historical data up to the current point for fitting
        train_data = data['Value'][:i]

        # Fit ARIMA model only if there are enough data points
        if len(train_data) > max(arima_order):
            try:
                arima_model = ARIMA(train_data, order=arima_order)
                arima_result = arima_model.fit()
                arima_pred = arima_result.forecast(steps=1).iloc[0]  # Use iloc to get the correct value
                arima_predictions.append(arima_pred)
            except (LinAlgError, ValueError) as e:
                arima_predictions.append(None)
                print(f"ARIMA model fitting failed at index {i} with error: {e}")
        else:
            arima_predictions.append(None)

        # Fit SARIMA model only if there are enough data points
        if len(train_data) > max(sarima_order[:3]) + sarima_order[3]:
            try:
                sarima_model = SARIMAX(train_data, order=sarima_order[:3], seasonal_order=sarima_order)
                sarima_result = sarima_model.fit(disp=False)
                sarima_pred = sarima_result.forecast(steps=1).iloc[0]  # Use iloc to get the correct value
                sarima_predictions.append(sarima_pred)
            except (LinAlgError, ValueError) as e:
                sarima_predictions.append(None)
                print(f"SARIMA model fitting failed at index {i} with error: {e}")
        else:
            sarima_predictions.append(None)

    # Add the predictions to the DataFrame, aligning them correctly
    data['ARIMA Predictions'] = [None] * window + arima_predictions
    data['SARIMA Predictions'] = [None] * window + sarima_predictions

    return data


# Time Series Relationships Analysis
#Verbesserungen: plot_series_with_predictors in plot_series Funktion verallgemeinern, die mehrere Zeitreihen in einem oder mehreren Diagrammen plottet und je nach Parameter separate Skalen anlegt oder alle einheitlich skaliert. Vorverarbeitung, Ausrichtung und Normalisierung in einzelne Methoden und in analyze_series_relationships und plot_series integrieren

# Preprocessing
def preprocess_series(series, max_consecutive_nans=3):
    """
    Preprocess the time series by interpolating small gaps and identifying the most recent valid segment.
    
    :param series: Time series (pandas Series)
    :param max_consecutive_nans: Maximum number of consecutive NaNs to fill by interpolation
    :return: Preprocessed time series (pandas Series)
    """
    # Convert 'nan' strings to actual NaN values
    series = series.replace('nan', np.nan).astype(float)

    # Fill small gaps with interpolation
    series_interpolated = series.interpolate(method='linear', limit=max_consecutive_nans, limit_direction='both')
    
    # Identify the most recent valid segment
    is_not_nan = ~series_interpolated.isna()
    seg_change = (is_not_nan != is_not_nan.shift()).cumsum()
    
    # List comprehension to get all valid segments
    valid_segments = [segment for _, segment in series_interpolated.groupby(seg_change) if segment.isna().sum() == 0]
    
    # If no valid segments found, return an empty series
    if not valid_segments:
        return pd.Series(dtype=series.dtype)
    
    # Return the most recent valid segment
    return valid_segments[-1]
def normalize_series(series, method='standard'):
    """
    Normalize the time series using the specified method.
    
    :param series: Time series (pandas Series)
    :param method: Normalization method ('standard' or 'minmax')
    :return: Normalized time series (pandas Series)
    """
    scaler = StandardScaler() if method == 'standard' else MinMaxScaler()
    series_scaled = scaler.fit_transform(series.values.reshape(-1, 1)).flatten()
    return pd.Series(series_scaled, index=series.index)
def check_stationarity(series, significance_level=0.05):
    """
    Check if the time series is stationary using the Augmented Dickey-Fuller test.
    
    :param series: Time series (pandas Series)
    :param significance_level: Significance level for the test
    :return: True if the series is stationary, False otherwise
    """
    result = adfuller(series)
    return result[1] < significance_level
def make_stationary(series):
    """
    Make the time series stationary by differencing and filling the first value with 0.
    
    :param series: Time series (pandas Series)
    :return: Stationary time series (pandas Series)
    """
    diff_series = series.diff().fillna(0)
    return diff_series
def ensure_stationarity(series, significance_level=0.05):
    # Check if the series are stationary
    if not check_stationarity(series, significance_level):
        series = make_stationary(series)
    return series
def extend_dataframe_dict(df, column_names, data_dict):
    for column_name in column_names:
        if column_name in data_dict:
            # Hole die Series aus dem Dictionary
            series = data_dict[column_name]
            
            # Filtere die Series, um nur die Einträge zu behalten, deren Index im DataFrame vorhanden ist
            filtered_series = series[series.index.isin(df.index)]
            
            # Füge die gefilterte Series als neue Spalte im DataFrame hinzu
            #df[column_name] = filtered_series
            #df.loc[:, column_name] = filtered_series
            df = pd.concat([df, filtered_series.rename(column_name)], axis=1)
    
    return df
def extend_dataframe(df, column_names, additional_df):
    for column_name in column_names:
        if column_name in additional_df.columns:
            # Hole die entsprechende Spalte aus dem zusätzlichen DataFrame
            series = additional_df[column_name]
            # Filtere die Series, um nur die Einträge zu behalten, deren Index im ursprünglichen DataFrame vorhanden ist
            filtered_series = series[series.index.isin(df.index)]
            # Füge die gefilterte Series als neue Spalte im ursprünglichen DataFrame hinzu
            #df[column_name] = filtered_series
            #df.loc[:, column_name] = filtered_series
            df = pd.concat([df, filtered_series.rename(column_name)], axis=1)
    return df
# Analysis
def analyze_series_relationship_deprecated(series1, series2, max_lag=24, plot=False):
    """
    Analyzes the relationship between two time series using cross-correlation and Granger causality tests.
    
    :param series1: First time series (pandas Series)
    :param series2: Second time series (pandas Series)
    :param max_lag: Maximum lag to consider for cross-correlation and Granger causality tests
    :param plot: If True, plot the cross-correlation
    :return: Dictionary with cross-correlation values and Granger causality test results
    """
    
    # Ensure the series are aligned and have the same length
    series1, series2 = series1.align(series2, join='inner')

    # Check if either series is constant
    if series1.nunique() == 1 or series2.nunique() == 1:
        print("One of the series is constant; skipping some tests.")
        return {
            "Comparison Period": (None, None),
            "Cross Correlation": np.array([]),
            "Granger Causality Test": {
                "P Values": {},
                "Best Lag": None
            }
        }
    
    # Ensure the series are stationary
    series1 = ensure_stationarity(series1)
    series2 = ensure_stationarity(series2)

    # Adjust the max_lag if there are insufficient observations
    max_lag = min(max_lag, (len(series1) - 2) // 3)
    print(f"Länge: {len(series1)}, Lag: {max_lag}")
    if max_lag < 1:
        print("Insufficient observations for Granger causality test.")
        return {
            "Comparison Period": (None, None),
            "Cross Correlation": np.array([]),
            "Granger Causality Test": {
                "P Values": {},
                "Best Lag": None
            }
        }
    
    # Calculate cross-correlation
    cross_corr = ccf(series1, series2)[:max_lag]
    
    if plot:
        # Plot cross-correlation
        plt.figure(figsize=(10, 5))
        plt.bar(range(max_lag), cross_corr)
        plt.xlabel('Lag')
        plt.ylabel('Cross-Correlation')
        plt.title('Cross-Correlation between Series1 and Series2')
        plt.show()
    
    # Perform Granger causality tests
    granger_results = grangercausalitytests(pd.concat([series1, series2], axis=1), max_lag)
    
    # Extract p-values from Granger causality tests and find the best lag
    p_values = {lag: test[0]['ssr_ftest'][1] for lag, test in granger_results.items()}
    best_lag = min(p_values, key=p_values.get)
    best_p_value = p_values[best_lag]
    
    # Pearson Correlation
    pearson_corr, pearson_p = pearsonr(series1, series2)
    
    # Spearman Rank Correlation
    spearman_corr, spearman_p = spearmanr(series1, series2)
    
    # Mutual Information
    mutual_info = mutual_info_regression(series1.values.reshape(-1, 1), series2.values)[0]
    
    # Dynamic Time Warping
    dtw_distance = dtw.distance(series1.values, series2.values)
    
    # Get the comparison period
    comparison_period = (series1.index.min(), series1.index.max())
    
    return {
        "Comparison Period": comparison_period,
        "Cross Correlation": cross_corr,
        "Granger Causality Test": {
            "P Values": p_values,
            "Best Lag": best_lag
        },
        "Pearson Correlation": {
            "Correlation": pearson_corr,
            "P Value": pearson_p
        },
        "Spearman Rank Correlation": {
            "Correlation": spearman_corr,
            "P Value": spearman_p
        },
        "Mutual Information": mutual_info,
        "Dynamic Time Warping": dtw_distance
    }
def compute_relationship_analysis(target_df, predictors_df, max_lag=12):
    # Stelle sicher, dass target_df und predictors_df die gleichen Indizes haben
    # target_df, predictors_df = target_df.align(predictors_df, join='inner', axis=0)

    # Überprüfe, ob die DataFrames nach der Interpolation noch Daten enthalten
    # if target_df.empty or predictors_df.empty:
    #     print("Keine gemeinsamen Datenpunkte nach Abgleich und NaN-Entfernung.")
    #     return pd.DataFrame()

    # Berechne die Kreuzkovarianzen für jeden Prädiktor
    results = []
    for pred_col in predictors_df.columns:
        for target_col in target_df.columns:
            target = target_df[target_col].copy()
            target = normalize_series(preprocess_series(target, max_consecutive_nans=3), method='standard')
            predictor = preprocess_series(predictors_df[pred_col], max_consecutive_nans=3)
            if predictor.empty:
                continue
            predictor = normalize_series(predictor, method='standard')

            if predictor.empty:
                continue
            
            target, predictor = target.align(predictor, join='inner')

            if len(target) == 0 or len(predictor) == 0:
                continue
            if len(target) != len(predictor):
                print(f"Skipping series {pred_col} due to length mismatch after alignment.")
                continue
            if target.nunique() == 1 or predictor.nunique() == 1:
                print(f"Skipping series {pred_col} due to being constant. If this series is not constant, check the target.")
                continue

            correlations = ccf(predictor, target, adjusted=False)[:max_lag]
            best_cr_lag = np.argmax(np.abs(correlations))
            max_correlation = correlations[best_cr_lag]

            covariances = ccovf(predictor, target, adjusted=False)[:max_lag]
            best_cv_lag = np.argmax(np.abs(covariances))
            max_covariance = covariances[best_cv_lag]

            target = ensure_stationarity(target)
            predictor = ensure_stationarity(predictor)
            maxlag = min(max_lag, (len(target) - 2) // 3)
            print(f"Länge: {len(target)}, Lag: {max_lag}")
            if max_lag < 1:
                print("Insufficient observations for Granger causality test.")
                continue
            try:
                granger_results = grangercausalitytests(pd.concat([target, predictor], axis=1), maxlag=maxlag)
                p_values = {lag: test[0]['ssr_ftest'][1] for lag, test in granger_results.items()}
                best_gc_lag = min(p_values, key=p_values.get)
                best_p_value = p_values[best_gc_lag]

                results.append({
                    'Target': target_col,
                    'Predictor': pred_col,
                    'Best Correlation': max_correlation,
                    'Best Correlation Lag': best_cr_lag,
                    'Correlations': correlations,
                    'Best Covariance': max_covariance,
                    'Best Covariance Lag': best_cv_lag,
                    #'Mean Absolute Covariance': np.abs(covariances).mean(),
                    'Covariances': covariances,
                    'Best Granger Causality P-Value': best_p_value,
                    'Best Granger Causality P-Value Lag': best_gc_lag,
                    'Granger Causality P-Values': p_values
                })
            except Exception as e:
                if type(e).__name__ == "InfeasibleTestError":
                    print(f"InfeasibleTestError catched, series {pred_col} is skipped:", e)
                else:
                    raise
          
    # Erstelle einen DataFrame aus den Ergebnissen
    results_df = pd.DataFrame(results)
    return results_df
# Predictors
def get_best_predictors_df_by_th(analysis_results_df, threshold=0.05):
    return analysis_results_df.reindex(analysis_results_df['Best Granger Causality P-Value'].sort_values(ascending=True).index)[analysis_results_df['Best Granger Causality P-Value'] < threshold][["Predictor", "Best Granger Causality P-Value"]]
def get_best_predictors_by_th(analysis_results_df, threshold=0.05):
    return get_best_predictors_df_by_th(analysis_results_df, threshold)["Predictor"].tolist()
def find_best_predictors(target_series, candidate_series_dict, n_best=10, max_lag=24, max_consecutive_nans=3, normalization_method='standard', return_all_results=False):
    """
    Finds the best predictors for a given target time series from a set of candidate time series.
    
    :param target_series: Target time series (pandas Series)
    :param candidate_series_dict: Dictionary of candidate time series (key: name, value: pandas Series)
    :param n_best: Number of best predictors to return (default is 10)
    :param max_lag: Maximum lag to consider for cross-correlation and Granger causality tests
    :param max_consecutive_nans: Maximum number of consecutive NaNs to fill by interpolation
    :param normalization_method: Normalization method ('standard' or 'minmax')
    :return: DataFrame with the n best predictors and their metrics
    """
    results = []
    
    # Normalize the target series
    target_series = normalize_series(preprocess_series(target_series, max_consecutive_nans), method=normalization_method)
    
    for name, series in candidate_series_dict.items():
        # Preprocess the candidate series
        series_preprocessed = preprocess_series(series, max_consecutive_nans)
        
        if series_preprocessed.empty:
            continue

        # Normalize the candidate series
        series_preprocessed = normalize_series(series_preprocessed, method=normalization_method)

        # Align series to find the longest overlapping time period
        aligned_target_series, aligned_series = target_series.align(series_preprocessed, join='inner')
        
        if len(aligned_target_series) == 0 or len(aligned_series) == 0:
            continue

        if len(aligned_target_series) != len(aligned_series):
            print(f"Skipping series {name} due to length mismatch after alignment.")
            continue
        
        print(name)
        metrics = analyze_series_relationship_deprecated(aligned_target_series, aligned_series, max_lag)
        if not metrics["Granger Causality Test"]["P Values"]:
            continue
        
        min_p_value = min(metrics["Granger Causality Test"]["P Values"].values())
        best_lag = metrics["Granger Causality Test"]["Best Lag"]
        comparison_period = metrics["Comparison Period"]
        results.append((name, min_p_value, metrics, best_lag, comparison_period))
    
    # Sort results by the minimum p-value
    results.sort(key=lambda x: x[1])

    best_results = results
    if not return_all_results:
        # Select the n best predictors
        best_results = results[:n_best]
    
    # Prepare the results DataFrame
    best_predictors = {
        'Predictor': [],
        'Min P-Value': [],
        'Best Lag': [],
        'Comparison Period': [],
        'Cross-Correlation': [],
        'Granger Causality P-Values': []
    }
    
    for result in best_results:
        best_predictors['Predictor'].append(result[0])
        best_predictors['Min P-Value'].append(result[1])
        best_predictors['Best Lag'].append(result[3])
        best_predictors['Comparison Period'].append(result[4])
        best_predictors['Cross-Correlation'].append(result[2]['Cross Correlation'])
        best_predictors['Granger Causality P-Values'].append(result[2]['Granger Causality Test']['P Values'])
    
    return pd.DataFrame(best_predictors)
#Application
def create_p_value_df(best_predictors_df, n_top_predictors=5):
    """
    Creates a DataFrame with the top predictors as columns and their p-values for each lag as rows.
    
    :param best_predictors_df: DataFrame with the best predictors and their metrics
    :param n_top_predictors: Number of top predictors to include in the DataFrame (default is 5)
    :return: DataFrame with p-values for each lag for the top predictors
    """
    # Sort the predictors by Min P-Value
    sorted_predictors = best_predictors_df.sort_values(by='Min P-Value').head(n_top_predictors)
    
    # Determine the maximum lag
    max_lag = max(
        max(p_values.keys()) for p_values in sorted_predictors['Granger Causality P-Values']
    )
    
    # Create a DataFrame to store the p-values for each lag
    p_value_df = pd.DataFrame(index=range(1, max_lag + 1))
    
    # Populate the DataFrame with p-values for each top predictor
    for _, row in sorted_predictors.iterrows():
        predictor_name = row['Predictor']
        p_values = row['Granger Causality P-Values']
        p_value_series = pd.Series(p_values, name=predictor_name)
        p_value_df = pd.concat([p_value_df, p_value_series], axis=1)
    
    return p_value_df
def find_best_predictors_oecd(n_best):
    oecd_df = pd.read_csv("../04 Data/OECD/OECD_integrated.csv")
    oecd_df.set_index("TIME_PERIOD", verify_integrity=True, inplace=True, drop=True)
    oecd_df.drop(columns=["Unnamed: 0"], inplace=True)
    oecd_df.index = pd.to_datetime(oecd_df.index)
    oecd_df.info()

    hg_oi_df = pd.read_csv("../04 Data/HOMAG/OrderIntake_series_monthly.csv")
    hg_oi_df.set_index("TIME_PERIOD", verify_integrity=True, inplace=True, drop=True)
    hg_oi_df.index = pd.to_datetime(hg_oi_df.index)
    hg_oi_df.info()

    target_series = hg_oi_df["T"]
    candidate_series_dict = {col: oecd_df[col] for col in oecd_df.columns}

    
    return target_series, candidate_series_dict, find_best_predictors(target_series, candidate_series_dict, n_best=n_best, return_all_results=True, max_lag=12)
# Plotting
def plot_series_with_predictors_plt(target_series, best_predictors_df, candidate_series_dict, n_best=10, n_per_plot=3, multiple_plots=False):
    """
    Plots the target series with the best predictors in a single plot or multiple plots using different scales.
    
    :param target_series: Target time series (pandas Series)
    :param best_predictors_df: DataFrame with the best predictors and their metrics
    :param candidate_series_dict: Dictionary of candidate time series (key: name, value: pandas Series)
    :param n_best: Number of best predictors to plot (default is 10)
    :param n_per_plot: Number of predictors per plot if multiple_plots is True (default is 3)
    :param multiple_plots: If True, create multiple plots with n_per_plot predictors each (default is False)
    """
    # Normalize the target series for consistent scaling
    scaler = StandardScaler()
    target_series_scaled = pd.Series(scaler.fit_transform(target_series.values.reshape(-1, 1)).flatten(), index=target_series.index)
    
    # Create a color map for the predictors
    cmap = plt.get_cmap('Greens')
    
    # If multiple_plots is True, split the best predictors into chunks of n_per_plot
    if multiple_plots:
        chunks = [best_predictors_df.iloc[i:i + n_per_plot] for i in range(0, n_best, n_per_plot)]
    else:
        chunks = [best_predictors_df.iloc[:n_best]]
    
    for chunk_index, chunk in enumerate(chunks):
        plt.figure(figsize=(14, 7))
        
        # Plot the target series
        plt.plot(target_series_scaled, label='Target', color='blue', linewidth=2)

        # Create a color map for the current chunk
        colors = [cmap(i / (len(chunk)+1)) for i in range(len(chunk)+1, 0, -1)]

        # Plot each predictor series with its own scale within the comparison period
        for i, row in enumerate(chunk.iterrows()):
            row = row[1]  # `row` is a tuple (index, Series), we need the Series
            name = row['Predictor']
            predictor_series = candidate_series_dict[name]
            
            # Extract the comparison period
            start_date, end_date = row['Comparison Period']
            
            # Filter the series to the comparison period
            predictor_series_filtered = predictor_series[start_date:end_date]
            
            # Normalize the predictor series
            scaler = StandardScaler()
            predictor_series_scaled = pd.Series(scaler.fit_transform(predictor_series_filtered.values.reshape(-1, 1)).flatten(), index=predictor_series_filtered.index)
            
            # Plot the predictor series
            plt.plot(predictor_series_scaled, label=f'{name} (Lag {row["Best Lag"]})', color=colors[i])
        
        if multiple_plots:
            plt.title(f'Target Series and Best Predictors (Plot {chunk_index + 1})')
        else:
            plt.title('Target Series and Best Predictors')
        plt.xlabel('Time')
        plt.ylabel('Normalized Values')
        plt.legend()
        plt.show()
def plot_time_series_with_moving_statistics(title="Order Intake", filepath="../04 Data/HOMAG/OrderIntake_series_monthly.csv", index_column="TIME_PERIOD", value_column="T"):
    hg_oi_df = pd.read_csv(filepath)

    hg_oi_statistics = calculate_moving_statistics(dates=hg_oi_df[index_column],values=hg_oi_df[value_column], window=4)

    plot_time_series(dates=hg_oi_statistics.index.tolist(),values=hg_oi_statistics[[col for col in hg_oi_statistics.columns if col != "Date"]],title=title)
def plot_time_series(dates, values, title='Time Series Plot', xlabel='Date', ylabel='Value'):
    """
    Plots a time series.

    :param dates: List of dates (datetime objects or strings in the format 'YYYY-MM-DD')
    :param values: List of values (float or int)
    :param title: Title of the plot (optional)
    :param xlabel: Label for the x-axis (optional)
    :param ylabel: Label for the y-axis (optional)
    """
    from datetime import datetime

    # Function to parse dates with multiple formats
    def parse_date(date_str):
        for fmt in ('%Y-%m-%d', '%Y-%m'):
            try:
                return datetime.strptime(date_str, fmt)
            except ValueError:
                continue
        raise ValueError(f"Date format for {date_str} not supported")

    # Convert dates to datetime objects if needed
    if isinstance(dates[0], str):
        dates = [parse_date(date) for date in dates]

    plt.figure(figsize=(20, 6))

    # Define a list of colors to use for the plots (excluding blue)
    colors = ['black', 'orange', 'green', 'red', 'purple', 'brown', 'pink', 'gray', 'olive', 'cyan', 'magenta']
    color_index = 0

    for col in values.columns:
        if col == "Value":
            plt.plot(dates, values[col], label= col, marker='o', linestyle='-', color='blue')
        else:
            plt.plot(dates, values[col], label= col, linestyle='--', color=colors[color_index % len(colors)])
            color_index += 1

    plt.title(title)
    plt.xlabel(xlabel)
    plt.ylabel(ylabel)
    plt.legend()

    # Format the x-axis for better readability
    plt.gca().xaxis.set_major_formatter(mdates.DateFormatter('%Y-%m-%d'))
    plt.gca().xaxis.set_major_locator(mdates.MonthLocator())
    plt.gcf().autofmt_xdate()  # Automatically rotate date labels

    plt.grid(True)
    plt.tight_layout()
    plt.show()
def plot_series_with_predictors(target_df, predictors_df, best_predictors_df, n_best=5, n_per_plot=3, multiple_plots=False, title="Order Intake", show=True):
    columns = best_predictors_df["Predictor"].head(n_best).values.tolist()
    target_and_predictors_df = extend_dataframe(target_df[["T"]], columns, additional_df=predictors_df)

    for column in columns:
        if column in target_and_predictors_df.columns:
            target_and_predictors_df[column] = normalize_series(preprocess_series(target_and_predictors_df[column], max_consecutive_nans=3), method='standard')
    target_and_predictors_df["T"] = normalize_series(preprocess_series(target_and_predictors_df["T"], max_consecutive_nans=3), method='standard')

    highlight_columns=[]#["Order Intake (total)"]

    # Erstellen der erweiterten Farbpalette
    k = 1
    total_columns = len(target_and_predictors_df.columns)
    greens_needed = total_columns - k
    if multiple_plots:
        greens_needed = n_per_plot
    # Nehme die ersten k Farben aus color_palette
    colors = color_palette[:k]
    colors.insert(0, color_palette[0])
    # Füge Grüntöne für die restlichen Spalten hinzu
    greens = plt.get_cmap('Greens', 2*greens_needed)
    for i in range(greens_needed):
        rgba = greens((i + 0.7*greens_needed) / (1.8*greens_needed))
        rgb = (int(rgba[0] * 255), int(rgba[1] * 255), int(rgba[2] * 255))
        colors.append(rgb)

    target_and_predictors_df = target_and_predictors_df.reset_index(drop=False).rename(columns={"TIME_PERIOD": "Date", "T": "Order Intake (total)"})

    if multiple_plots:
        # Split predictors into chunks of size n_per_plot
        num_plots = (len(columns) + n_per_plot - 1) // n_per_plot  # Ceiling division
        for i in range(num_plots):
            start_index = i * n_per_plot
            end_index = min((i + 1) * n_per_plot, len(columns))
            plot_columns = ["Order Intake (total)"] + columns[start_index:end_index]

            # Create a subset DataFrame for the current plot
            subset_df = target_and_predictors_df[["Date"] + plot_columns]

            plot, plot_imgstream = create_plot(
                df=subset_df,
                highlight_columns=highlight_columns,
                color_palette=colors,
                chart_title=f'{title} - Plot {i + 1}',
                xlabel='Datum',
                ylabel='Normalisierte Werte',
                #scale_factor=1e-6,
                default_columns=subset_df.columns
            )
            plot.show()
    else:
        plot, plot_imgstream = create_plot(
            df=target_and_predictors_df,
            highlight_columns=highlight_columns,
            color_palette=colors,
            chart_title=title,
            xlabel='Datum',
            ylabel='Normalisierte Werte',
            #scale_factor=1e-6,
            default_columns=["Order Intake (total)"]
        )
        if show:
            plot.show()
        return plot, plot_imgstream
def create_plot(df, highlight_columns, chart_title, xlabel, ylabel, color_palette, scale_factor=None, default_columns=None):
    # DataFrame skalieren, wenn ein Skalierungsfaktor angegeben wurde
    if scale_factor is not None:
        numeric_columns = df.select_dtypes(include=[np.number])
        df[numeric_columns.columns] = numeric_columns * scale_factor

    # Farben vorbereiten
    colors = {col: f'rgb{color_palette[i % len(color_palette)]}' for i, col in enumerate(df.columns) if col != 'Date'}

    # Daten in das lange Format umwandeln
    df_long = df.melt(id_vars='Date', var_name='Variable', value_name='Value')

    # Plotly Express verwenden, um das Diagramm zu erstellen
    fig = px.line(df_long, x='Date', y='Value', color='Variable', title=chart_title,
                  labels={'Value': ylabel, 'Date': xlabel}, color_discrete_map=colors)

    # Anpassung der Linienbreite für hervorgehobene Spalten
    for trace in fig.data:
        trace.line.width = 2.5 if trace.name in highlight_columns else 1.5

    # Standardmäßig angezeigte Spalten festlegen und Dropdown-Menü hinzufügen, wenn default_columns gesetzt ist
    if default_columns is not None:
        for trace in fig.data:
            trace.visible = trace.name in default_columns

        # Dropdown-Menü für die Auswahl von Spalten hinzufügen
        buttons = []
        buttons.append(dict(label='Alle', method='update', args=[{'visible': [True] * len(fig.data)},
                                                                 {'title': chart_title}]))
        for column in df.columns:
            if column != 'Date':
                buttons.append(dict(
                    label=column,
                    method='update',
                    args=[{'visible': [trace.name == column or trace.name in default_columns for trace in fig.data]},
                          {'title': f'{chart_title} - {column}'}]
                ))
        

        fig.update_layout(
            updatemenus=[dict(
                active=1,
                buttons=buttons,
                direction='down',
                showactive=True
            )]
        )

    fig.update_xaxes(
        rangeslider_visible= True,
        rangeselector=dict(
            buttons=list([
                #dict(count=1, label="1m", step="month", stepmode="backward"),
                dict(count=6, label="6m", step="month", stepmode="backward"),
                dict(count=1, label="YTD", step="year", stepmode="todate"),
                dict(count=1, label="1y", step="year", stepmode="backward"),
                dict(count=5, label="5y", step="year", stepmode="backward"),
                dict(step="all")
            ])
        )
    )

    # Diagrammhöhe anpassen
    fig.update_layout(height=700)  # Erhöhen Sie die Höhe des Diagramms

    # Diagramm in ein BytesIO-Objekt speichern
    img_stream = imgstream_pyplot(fig)

    return fig, img_stream
def imgstream_pyplot(fig):
    img_stream = BytesIO()
    fig.write_image(img_stream, format='png')
    img_stream.seek(0)
    return img_stream
# Documentation
color_palette = [
    (0, 25, 65),     # Dunkelblau
    (0, 160, 220),   # Hellblau
    (255, 200, 0),   # Gelb
    (122, 122, 122), # Dunkelgrau
    (163, 163, 163), # Mittelgrau
    (204, 204, 204), # Hellgrau
    (140, 210, 30),  # Grün
    (240, 50, 50),   # Rot
    (255, 130, 0)    # Orange
]
def create_plot_plt(df, highlight_columns, chart_title, xlabel, ylabel, color_palette, scale_factor=None):
    # DataFrame skalieren, wenn ein Skalierungsfaktor angegeben wurde
    if scale_factor is not None:
        numeric_columns = df.select_dtypes(include=[np.number])
        df[numeric_columns.columns] = numeric_columns * scale_factor

    # Matplotlib Diagramm erstellen
    plt.figure(figsize=(10, 6))
    color_index = 0
    for column in df.columns:
        if column != 'Date':
            color = '#%02x%02x%02x' % color_palette[color_index % len(color_palette)]
            plt.plot(df['Date'], df[column], label=column, color=color, linewidth=2.5 if column in highlight_columns else 1.5)
            color_index += 1

    plt.xlabel(xlabel)
    plt.ylabel(ylabel)
    plt.title(chart_title)
    plt.legend()
    plt.tight_layout()

    # Diagramm in ein BytesIO-Objekt speichern und zurückgeben
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight', dpi=300)
    plt.close()
    img_stream.seek(0)
    return img_stream
def add_slide_with_image(pptx_path, title, img_stream):
    # Präsentation laden
    prs = Presentation(pptx_path)
    
    # Neue Folie hinzufügen
    slide_layout = prs.slide_layouts[4]  # Verwenden eines Layouts mit Platzhaltern
    slide = prs.slides.add_slide(slide_layout)

    # Folientitel hinzufügen
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Berechnen der Bildgröße in Inches
    image = Image.open(img_stream)
    width, height = image.size
    img_width_inch = width / 300  # Umrechnung von Pixel in Inches (bei 300 dpi)
    img_height_inch = height / 300

    # Versuchen, einen geeigneten Platzhalter für das Bild zu finden
    placeholder_found = False
    for shape in slide.shapes:
        if shape.is_placeholder and shape.has_text_frame:
            # Prüfen, ob der Platzhalter ein Textfeld ist
            if not shape.text_frame.text:
                # Leeres Textfeld gefunden, Bild hier einfügen
                placeholder_width = shape.width / 914400  # Umrechnung von EMU in Inches
                placeholder_height = shape.height / 914400

                # Berechnen des Skalierungsfaktors
                scale_factor = min(placeholder_width / img_width_inch, placeholder_height / img_height_inch)

                # Angepasste Bildgröße
                new_width = img_width_inch * scale_factor
                new_height = img_height_inch * scale_factor

                # Berechnen der Zentrierungsposition innerhalb des Platzhalters
                placeholder_left = shape.left + (shape.width - Inches(new_width)) / 2
                placeholder_top = shape.top + (shape.height - Inches(new_height)) / 2

                # Bild im Platzhalter zentriert einfügen
                slide.shapes.add_picture(img_stream, placeholder_left, placeholder_top, width=Inches(new_width), height=Inches(new_height))
                placeholder_found = True
                break

    # Wenn kein geeigneter Platzhalter gefunden wurde, das Bild auf der Folie zentriert einfügen
    if not placeholder_found:
        left = (prs.slide_width - Inches(img_width_inch)) / 2
        top = (prs.slide_height - Inches(img_height_inch)) / 2
        slide.shapes.add_picture(img_stream, left, top, width=Inches(img_width_inch), height=Inches(img_height_inch))

    # Präsentation speichern
    prs.save(pptx_path)